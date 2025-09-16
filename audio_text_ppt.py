import os
import whisper
import warnings
from tkinter import Tk
from tkinter.filedialog import askopenfilename
from pptx import Presentation

warnings.filterwarnings("ignore", message="FP16 is not supported on CPU")

def audio_to_text(audio_path, max_chars_per_slide=400):
    try:
        print(" Generating transcript... (this may take some time)")
        model = whisper.load_model("base", device="cpu")

        # Transcribe audio
        result = model.transcribe(audio_path)
        transcript = result["text"].strip()

        # Save transcript to TXT
        txt_file = os.path.splitext(audio_path)[0] + ".txt"
        with open(txt_file, "w", encoding="utf-8") as f:
            f.write(transcript)
        print(f"Transcript saved: {txt_file}")

        # Create PPTX from transcript
        pptx_file = os.path.splitext(audio_path)[0] + ".pptx"
        prs = Presentation()

        # Split transcript into chunks (max characters per slide)
        chunks = [transcript[i:i+max_chars_per_slide] for i in range(0, len(transcript), max_chars_per_slide)]

        for i, chunk in enumerate(chunks, start=1):
            slide_layout = prs.slide_layouts[1]  # Title & Content layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Transcript Part {i}"
            slide.placeholders[1].text = chunk

        prs.save(pptx_file)
        print(f"PPTX created: {pptx_file}")

    except Exception as e:
        print(f"Error: {str(e)}")

if __name__ == "__main__":
    Tk().withdraw()
    audio_file = askopenfilename(
        title="Select an audio file",
        filetypes=[("Audio files", "*.mp3 *.wav *.m4a *.flac"), ("All files", "*.*")]
    )

    if audio_file:
        audio_to_text(audio_file, max_chars_per_slide=400)
    else:
        print("No file selected.")
