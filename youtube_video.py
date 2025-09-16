import yt_dlp
import moviepy as mp
import os
import whisper
import warnings
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

warnings.filterwarnings("ignore", message="FP16 is not supported on CPU")

def download_and_convert(url, save_path=".", words_per_bullet=15, bullets_per_slide=5):
    try:
        # Step 1: Download video
        ydl_opts = {"outtmpl": f"{save_path}/%(title)s.%(ext)s"}
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            video_file = ydl.prepare_filename(info)
        print(f"Download complete: {video_file}")

        # Step 2: Convert to MP3
        mp3_file = os.path.splitext(video_file)[0] + ".mp3"
        video = mp.VideoFileClip(video_file)
        audio = video.audio
        audio.write_audiofile(mp3_file)
        audio.close()
        video.close()
        print(f"MP3 saved: {mp3_file}")

        # Step 3: Transcribe audio
        print("📝 Generating transcript... (this may take some time)")
        model = whisper.load_model("base", device="cpu")
        result = model.transcribe(mp3_file)
        transcript = result["text"].strip()

        # Step 4: Save transcript to TXT
        txt_file = os.path.splitext(video_file)[0] + ".txt"
        with open(txt_file, "w", encoding="utf-8") as f:
            f.write(transcript)
        print(f"Transcript saved: {txt_file}")

        # Step 5: Create PPTX
        pptx_file = os.path.splitext(video_file)[0] + ".pptx"
        prs = Presentation()

        # Split transcript into bullets
        words = transcript.split()
        bullets = [" ".join(words[i:i+words_per_bullet]) for i in range(0, len(words), words_per_bullet)]

        # Group bullets into slides
        for i in range(0, len(bullets), bullets_per_slide):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Slide background color
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 224)  # Light yellow

            # Title box
            left, top, width, height = Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = f"Transcript Part {(i // bullets_per_slide) + 1}"
            p = title_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.font.name = "Mangal"

            # Content box
            left, top, width, height = Inches(0.5), Inches(1), Inches(9), Inches(5.5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            tf = content_box.text_frame
            tf.clear()
            for bullet in bullets[i:i+bullets_per_slide]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                font = p.font
                font.name = "Mangal"
                font.size = Pt(18)
                font.color.rgb = RGBColor(0, 0, 0)

            # Slide number
            left, top, width, height = Inches(8), Inches(6.8), Inches(1), Inches(0.3)
            sn_box = slide.shapes.add_textbox(left, top, width, height)
            sn_frame = sn_box.text_frame
            sn_frame.text = f"{(i // bullets_per_slide) + 1}"
            sn_frame.paragraphs[0].font.size = Pt(14)
            sn_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            sn_frame.paragraphs[0].font.name = "Mangal"

            # Fade transition
            slide_elm = slide._element
            transition_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" transition="fade" speed="med"/>'
            slide_elm.insert(2, parse_xml(transition_xml))

        prs.save(pptx_file)
        print(f"PPTX created: {pptx_file}")

        # Auto-open PPTX
        try:
            os.startfile(pptx_file)  # Windows only
            print("Opening PPTX file...")
        except Exception as open_err:
            print(f"⚠ Could not auto-open PPTX: {open_err}")

    except Exception as e:
        print(f"Error: {str(e)}")

if __name__ == "__main__":
    url = input("Enter YouTube video URL: ")
    download_and_convert(url)
